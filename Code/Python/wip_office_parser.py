# pip install python-docx pdfplumber openpyxl xlrd python-pptx odfpy
import os
from docx import Document
import pdfplumber
import openpyxl
import xlrd
from pptx import Presentation
from odf import text, teletype, spreadsheet, presentation

# Load keywords from a file
def load_keywords(filepath):
    with open(filepath, 'r') as file:
        return [line.strip() for line in file.readlines()]

# Reading functions for various file types
def read_txt_file(filepath):
    with open(filepath, 'r') as file:
        return file.read()

def read_docx_file(filepath):
    doc = Document(filepath)
    return '\n'.join([para.text for para in doc.paragraphs])

def read_pdf_file(filepath):
    text = ''
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ''
    return text

def read_xlsx_file(filepath):
    workbook = openpyxl.load_workbook(filepath)
    sheet = workbook.active
    return '\n'.join([' '.join([str(cell) for cell in row if cell is not None]) for row in sheet.iter_rows(values_only=True)])

def read_xls_file(filepath):
    workbook = xlrd.open_workbook(filepath)
    sheet = workbook.sheet_by_index(0)
    return '\n'.join([' '.join([str(sheet.cell_value(row, col)) for col in range(sheet.ncols)]) for row in range(sheet.nrows)])

def read_pptx_file(filepath):
    prs = Presentation(filepath)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def read_odt_file(filepath):
    odtFile = text.load(filepath)
    allText = teletype.extractText(odtFile.documentElement)
    return allText

def read_ods_file(filepath):
    odsFile = spreadsheet.load(filepath)
    text = ''
    for table in odsFile.spreadsheet.getElementsByType(spreadsheet.Table):
        for row in table.getElementsByType(spreadsheet.TableRow):
            for cell in row.getElementsByType(spreadsheet.TableCell):
                repeat = cell.getAttribute("numbercolumnsrepeated")
                if not repeat: repeat = 1
                for _ in range(int(repeat)):  # Handle repeated cells
                    if cell.firstChild:
                        text += teletype.extractText(cell) + ' '
            text += '\n'
    return text

def read_odp_file(filepath):
    odpFile = presentation.load(filepath)
    text = ''
    for page in odpFile.getElementsByType(presentation.Page):
        for frame in page.getElementsByType(presentation.Frame):
            for textbox in frame.getElementsByType(text.P):
                text += teletype.extractText(textbox) + '\n'
    return text

# Search for keywords in the provided text
def search_for_keywords(text, keywords, file):
    for keyword in keywords:
        if keyword.lower() in text.lower():
            print(f"KEYWORD FOUND: {keyword}, FILE: {file}")

# Process each file based on its type
def process_file(filepath, keywords):
    text = ''
    if filepath.endswith('.txt'):
        text = read_txt_file(filepath)
    elif filepath.endswith('.docx'):
        text = read_docx_file(filepath)
    elif filepath.endswith('.pdf'):
        text = read_pdf_file(filepath)
    elif filepath.endswith('.xlsx'):
        text = read_xlsx_file(filepath)
    elif filepath.endswith('.xls'):
        text = read_xls_file(filepath)
    elif filepath.endswith('.pptx'):
        text = read_pptx_file(filepath)
    elif filepath.endswith('.odt'):
        text = read_odt_file(filepath)
    elif filepath.endswith('.ods'):
        text = read_ods_file(filepath)
    elif filepath.endswith('.odp'):
        text = read_odp_file(filepath)
    else:
        print(f"Unsupported file type: {filepath}")
        return
    search_for_keywords(text, keywords, filepath)

# Walk through the folder and process each file
def walk_through_folder(folder, keywords):
    for root, dirs, files in os.walk(folder):
        for file in files:
            filepath = os.path.join(root, file)
            process_file(filepath, keywords)

keywords_file = 'keywords.txt'
folder = 'path_to_your_folder'
keywords = load_keywords(keywords_file)

walk_through_folder(folder, keywords)
