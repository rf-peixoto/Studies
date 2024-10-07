import os
import shutil
import logging
import yaml
import sys
import argparse
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import defaultdict
import threading
import json
import re
import email
from email import policy
from email.parser import BytesParser

# Import libraries for different file types
import docx
import pptx
import csv
import json as json_module

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None  # Optional, for PDF support

try:
    from odf import text, teletype
    from odf.opendocument import load as load_odf
except ImportError:
    load_odf = None  # Optional, for LibreOffice support

try:
    import openpyxl
except ImportError:
    openpyxl = None  # Optional, for .xlsx support

try:
    import toml
except ImportError:
    toml = None  # Optional, for TOML support

try:
    import extract_msg
except ImportError:
    extract_msg = None  # Optional, for .msg support

import xml.etree.ElementTree as ET  # Standard library for XML parsing


# Thread-safe statistics collector
class Statistics:
    def __init__(self):
        self.lock = threading.Lock()
        self.keyword_group_counts = defaultdict(int)
        self.keyword_counts = defaultdict(lambda: defaultdict(int))
        self.filetype_counts = defaultdict(int)
        self.files_with_multiple_keywords = 0

    def update_keyword_group(self, group):
        with self.lock:
            self.keyword_group_counts[group] += 1

    def update_keyword(self, group, keyword):
        with self.lock:
            self.keyword_counts[group][keyword] += 1

    def update_filetype(self, filetype):
        with self.lock:
            self.filetype_counts[filetype] += 1

    def increment_multiple_keywords(self):
        with self.lock:
            self.files_with_multiple_keywords += 1

    def to_dict(self):
        with self.lock:
            return {
                "Keyword_Groups": {k: v for k, v in self.keyword_group_counts.items()},
                "Keywords": {k: dict(v) for k, v in self.keyword_counts.items()},
                "FileTypes": {k: v for k, v in self.filetype_counts.items()},
                "Files_With_Multiple_Keywords": self.files_with_multiple_keywords
            }


# Function to load and validate configuration
def load_config(config_path):
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    if 'Keyword_Groups' not in config:
        logging.error("Configuration file is missing 'Keyword_Groups' section.")
        sys.exit(1)
    keyword_groups = config['Keyword_Groups']
    if not isinstance(keyword_groups, dict):
        logging.error("'Keyword_Groups' should be a dictionary.")
        sys.exit(1)
    # Validate that each group has a list of keywords
    for group, keywords in keyword_groups.items():
        if not isinstance(keywords, list):
            logging.error(f"Keywords for group '{group}' should be a list.")
            sys.exit(1)
    return keyword_groups


# Function to setup logging
def setup_logging(log_file):
    logging.basicConfig(
        filename=log_file,
        filemode='w',
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )
    # Also log to console
    console = logging.StreamHandler()
    console.setLevel(logging.INFO)
    formatter = logging.Formatter('%(levelname)s - %(message)s')
    console.setFormatter(formatter)
    logging.getLogger().addHandler(console)


# Function to extract text from different file types
def extract_text(file_path):
    ext = file_path.suffix.lower()
    text_content = ""
    try:
        if ext in ['.txt', '.csv', '.json', '.sql', '.conf', '.cfg']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        elif ext in ['.docx', '.docm', '.dotx', '.dotm']:
            doc = docx.Document(file_path)
            text_content = '\n'.join([para.text for para in doc.paragraphs])
        elif ext in ['.xlsx', '.xls', '.xlsm', '.xltx', '.xltm']:
            if ext == '.xlsx' and openpyxl:
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        row_text = ' '.join([str(cell) for cell in row if cell is not None])
                        text_content += row_text + '\n'
                wb.close()
            else:
                # For .xls and other formats handled by xlrd
                try:
                    import xlrd
                    workbook = xlrd.open_workbook(file_path, on_demand=True)
                    sheets = workbook.sheet_names()
                    for sheet in sheets:
                        worksheet = workbook.sheet_by_name(sheet)
                        for row in range(worksheet.nrows):
                            text_content += ' '.join([str(cell) for cell in worksheet.row(row)]) + '\n'
                except ImportError:
                    logging.error("xlrd is not installed. Install it or switch to openpyxl for .xlsx files.")
        elif ext in ['.pptx', '.pptm', '.potx', '.potm']:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + '\n'
        elif ext in ['.odt', '.ods', '.odp']:
            if load_odf:
                odf_doc = load_odf(file_path)
                if ext == '.odt':
                    allparas = odf_doc.getElementsByType(text.P)
                    text_content = '\n'.join([teletype.extractText(p) for p in allparas])
                elif ext == '.ods':
                    sheets = odf_doc.spreadsheet.getElementsByType(text.Table)
                    for sheet in sheets:
                        rows = sheet.getElementsByType(text.TableRow)
                        for row in rows:
                            cells = row.getElementsByType(text.TableCell)
                            row_text = ' '.join([teletype.extractText(cell) for cell in cells])
                            text_content += row_text + '\n'
                elif ext == '.odp':
                    slides = odf_doc.getElementsByType(text.Slide)
                    for slide in slides:
                        allparas = slide.getElementsByType(text.P)
                        text_content += '\n'.join([teletype.extractText(p) for p in allparas]) + '\n'
            else:
                logging.warning(f"LibreOffice support is not available. Install 'odfpy' to handle {ext} files.")
        elif ext == '.pdf' and PyPDF2:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text_content += extracted + '\n'
        elif ext == '.xml':
            try:
                tree = ET.parse(file_path)
                root = tree.getroot()
                text_content = extract_text_from_xml(root)
            except ET.ParseError as e:
                logging.error(f"Error parsing XML file {file_path}: {e}")
            except Exception as e:
                logging.error(f"Unexpected error processing XML file {file_path}: {e}")
        elif ext == '.eml':
            try:
                with open(file_path, 'rb') as f:
                    msg = BytesParser(policy=policy.default).parse(f)
                # Extract plain text parts
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        text_content += part.get_content()
                    elif part.get_content_type() == 'text/html':
                        # Optionally, you can strip HTML tags or use BeautifulSoup to extract text
                        html_content = part.get_content()
                        # Simple regex to remove HTML tags
                        clean_text = re.sub('<[^<]+?>', '', html_content)
                        text_content += clean_text + '\n'
            except Exception as e:
                logging.error(f"Error reading EML file {file_path}: {e}")
        elif ext == '.msg':
            if extract_msg:
                try:
                    msg = extract_msg.Message(file_path)
                    msg_sender = msg.sender
                    msg_date = msg.date
                    msg_subject = msg.subject
                    msg_body = msg.body
                    msg_text = f"Sender: {msg_sender}\nDate: {msg_date}\nSubject: {msg_subject}\n\n{msg_body}"
                    text_content = msg_text
                except Exception as e:
                    logging.error(f"Error reading MSG file {file_path}: {e}")
            else:
                logging.error("extract-msg library is not installed. Install it to handle .msg files.")
        elif ext in ['.ini', '.toml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        else:
            logging.warning(f"Unsupported file type: {file_path}")
    except Exception as e:
        logging.error(f"Error reading {file_path}: {e}")
    return text_content.lower()


def extract_text_from_xml(element):
    """Recursively extract text from XML elements."""
    text_content = element.text or ""
    for child in element:
        text_content += '\n' + extract_text_from_xml(child)
        if child.tail:
            text_content += '\n' + child.tail
    return text_content


# Function to compile regex patterns for all keywords
def compile_keyword_patterns(keyword_groups):
    compiled_patterns = {}
    for group, keywords in keyword_groups.items():
        compiled_patterns[group] = []
        for keyword in keywords:
            # Escape regex special characters in keyword
            escaped_keyword = re.escape(keyword.lower())
            # Determine if keyword is a single word or a phrase
            if ' ' in keyword:
                # For phrases, use word boundaries around the entire phrase
                pattern = r'(?<!\w)' + escaped_keyword + r'(?!\w)'
            else:
                # For single words, use word boundaries
                pattern = r'\b' + escaped_keyword + r'\b'
            try:
                compiled = re.compile(pattern)
                compiled_patterns[group].append((keyword, compiled))
            except re.error as e:
                logging.error(f"Invalid regex pattern for keyword '{keyword}' in group '{group}': {e}")
    return compiled_patterns


# Function to search keywords in text using compiled regex patterns
def search_keywords(text_content, compiled_patterns):
    found_groups = {}
    for group, patterns in compiled_patterns.items():
        for keyword, pattern in patterns:
            if pattern.search(text_content):
                if group not in found_groups:
                    found_groups[group] = set()
                found_groups[group].add(keyword)
    return found_groups


# Function to handle file copying with naming conflict resolution
def copy_file_to_groups(file_path, groups, output_dir):
    copied_groups = []
    for group in groups:
        group_folder = output_dir / group
        group_folder.mkdir(parents=True, exist_ok=True)
        destination = group_folder / file_path.name
        # Handle naming conflicts
        if destination.exists():
            base, ext = os.path.splitext(file_path.name)
            counter = 1
            while True:
                new_name = f"{base}_{counter}{ext}"
                new_destination = group_folder / new_name
                if not new_destination.exists():
                    destination = new_destination
                    break
                counter += 1
        try:
            shutil.copy2(file_path, destination)
            logging.info(f"Copied {file_path} to {destination}")
            copied_groups.append(group)
        except Exception as e:
            logging.error(f"Error copying {file_path} to {destination}: {e}")
    return copied_groups


# Function to process a single file
def process_file(file_path, compiled_patterns, output_dir, stats):
    logging.info(f"Processing file: {file_path}")
    text_content = extract_text(file_path)
    if text_content:
        found = search_keywords(text_content, compiled_patterns)
        if found:
            # Update statistics
            for group, keywords in found.items():
                stats.update_keyword_group(group)
                for keyword in keywords:
                    stats.update_keyword(group, keyword)
            if len(found) > 1:
                stats.increment_multiple_keywords()
            # Copy file to corresponding group folders
            copied_groups = copy_file_to_groups(file_path, found.keys(), output_dir)
            # Log found keywords
            for group, keywords in found.items():
                for keyword in keywords:
                    logging.info(f"Found keyword '{keyword}' in {file_path} for group '{group}'")
        # Update filetype statistics
        stats.update_filetype(file_path.suffix.lower())
    else:
        logging.warning(f"No text extracted from {file_path}")


# Function to generate statistics output
def generate_statistics(stats, output_dir):
    stats_data = stats.to_dict()
    stats_file = output_dir / "statistics.json"
    try:
        with open(stats_file, 'w', encoding='utf-8') as f:
            json.dump(stats_data, f, indent=4)
        logging.info(f"Statistics written to {stats_file}")
    except Exception as e:
        logging.error(f"Error writing statistics to {stats_file}: {e}")


# Main function
def main(data_dir, config_path, output_dir, log_file):
    setup_logging(log_file)
    keyword_groups = load_config(config_path)
    compiled_patterns = compile_keyword_patterns(keyword_groups)
    stats = Statistics()
    data_path = Path(data_dir)
    if not data_path.exists():
        logging.error(f"Data directory does not exist: {data_dir}")
        sys.exit(1)
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    # Supported file extensions
    supported_exts = [
        '.txt', '.csv', '.json', '.sql', '.conf', '.cfg',
        '.docx', '.docm', '.dotx', '.dotm',
        '.xlsx', '.xls', '.xlsm', '.xltx', '.xltm',
        '.pptx', '.pptm', '.potx', '.potm',
        '.odt', '.ods', '.odp',
        '.pdf', '.xml', '.ini', '.toml',
        '.eml', '.msg'  # Added email formats
    ]
    # Gather all supported files
    all_files = []
    for root, dirs, files in os.walk(data_path):
        for file in files:
            file_path = Path(root) / file
            if file_path.suffix.lower() in supported_exts:
                all_files.append(file_path)
            else:
                logging.warning(f"Skipping unsupported file type: {file_path}")

    if not all_files:
        logging.info("No supported files found to process.")
        return

    # Define number of threads (you can adjust this number based on your system)
    max_workers = min(32, os.cpu_count() + 4)

    # Process files in parallel
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [
            executor.submit(process_file, file_path, compiled_patterns, output_path, stats)
            for file_path in all_files
        ]
        for future in as_completed(futures):
            try:
                future.result()
            except Exception as e:
                logging.error(f"Unhandled exception during file processing: {e}")

    # Generate statistics output
    generate_statistics(stats, output_path)

    logging.info("Processing completed.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Keyword Search and File Organizer")
    parser.add_argument('--data', default='./data', help="Path to the data directory (default: ./data)")
    parser.add_argument('--config', default='./config.yaml', help="Path to the config.yaml file (default: ./config.yaml)")
    parser.add_argument('--output', default='./output', help="Directory to store copied files (default: ./output)")
    parser.add_argument('--log', default='./script.log', help="Log file path (default: ./script.log)")
    args = parser.parse_args()
    main(args.data, args.config, args.output, args.log)
